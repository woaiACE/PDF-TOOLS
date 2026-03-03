import os
from pypdf import PdfWriter, PdfReader

def merge_pdfs(pdf_list, output_path):
    merger = PdfWriter()
    try:
        for pdf in pdf_list:
            merger.append(pdf)
        with open(output_path, "wb") as f:
            merger.write(f)
        return True, "合并成功！"
    except Exception as e:
        return False, f"合并失败: {str(e)}"
    finally:
        merger.close()

def parse_page_ranges(range_str, total_pages):
    pages = set()
    parts = [p.strip() for p in range_str.split(',') if p.strip()]
    for part in parts:
        if '-' in part:
            try:
                start, end = map(int, part.split('-'))
                if start > end:
                    start, end = end, start
                for p in range(start, end + 1):
                    if 1 <= p <= total_pages:
                        pages.add(p - 1)
            except ValueError:
                raise ValueError(f"无效的页码范围格式: {part}")
        else:
            try:
                p = int(part)
                if 1 <= p <= total_pages:
                    pages.add(p - 1)
            except ValueError:
                raise ValueError(f"无效的页码格式: {part}")
    return sorted(list(pages))

def split_pdf_extract(pdf_path, output_path, range_str):
    try:
        reader = PdfReader(pdf_path)
        total_pages = len(reader.pages)
        pages_to_extract = parse_page_ranges(range_str, total_pages)
        if not pages_to_extract:
            return False, "没有选中任何有效的页面。"
        writer = PdfWriter()
        for p in pages_to_extract:
            writer.add_page(reader.pages[p])
        with open(output_path, "wb") as f:
            writer.write(f)
        return True, "提取成功！"
    except Exception as e:
        return False, f"提取失败: {str(e)}"

def split_pdf_all(pdf_path, output_dir):
    try:
        reader = PdfReader(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        for i, page in enumerate(reader.pages):
            writer = PdfWriter()
            writer.add_page(page)
            output_path = os.path.join(output_dir, f"{base_name}_page_{i+1}.pdf")
            with open(output_path, "wb") as f:
                writer.write(f)
        return True, "拆分成功！"
    except Exception as e:
        return False, f"拆分失败: {str(e)}"

import fitz # PyMuPDF
from pdf2docx import Converter
from PIL import Image

def compress_pdf(input_path, output_path):
    """
    使用 PyMuPDF 进行基础的垃圾回收和图像压缩优化
    """
    try:
        doc = fitz.open(input_path)
        doc.save(
            output_path,
            garbage=4,          # 最大的垃圾回收力度
            deflate=True,       # 压缩未压缩的流
            clean=True          # 净化和清理内容流
        )
        doc.close()
        return True, "PDF压缩成功！"
    except Exception as e:
        return False, f"压缩失败: {str(e)}"

def pdf_to_word(input_path, output_path):
    """
    使用 pdf2docx 库将 PDF 转换为 Word (docx)
    """
    try:
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        return True, "PDF成功转换为Word文档！"
    except Exception as e:
        return False, f"转换失败: {str(e)}"

def pdf_to_excel(input_path, output_path):
    """
    使用 pdfplumber 提取表格并用 pandas 保存到 Excel
    (简单的表格提取逻辑，复杂的可能需要专业工具)
    """
    try:
        import pdfplumber
        import pandas as pd

        all_tables = []
        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    df = pd.DataFrame(table[1:], columns=table[0])
                    all_tables.append(df)

        if not all_tables:
            return False, "未能从PDF中检测到可提取的表格数据。"

        with pd.ExcelWriter(output_path) as writer:
            for i, df in enumerate(all_tables):
                df.to_excel(writer, sheet_name=f'Sheet{i+1}', index=False)
        return True, "PDF表格成功提取至Excel！"
    except Exception as e:
        return False, f"转换失败: {str(e)}"

def pdf_to_ppt(input_path, output_path):
    """
    将 PDF 每页作为图片插入到 PPT 幻灯片中。
    （保留了外观，但不是可编辑的文本）
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
        import tempfile

        prs = Presentation()
        # default slide size is 4:3 (10 x 7.5 inches)
        # we can adjust it based on the first pdf page aspect ratio if needed
        blank_slide_layout = prs.slide_layouts[6]

        doc = fitz.open(input_path)

        with tempfile.TemporaryDirectory() as tmpdirname:
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=150)
                img_path = os.path.join(tmpdirname, f"page_{page_num}.png")
                pix.save(img_path)

                slide = prs.slides.add_slide(blank_slide_layout)

                # Fit image to slide (very basic)
                slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

        prs.save(output_path)
        doc.close()
        return True, "PDF成功转换为PowerPoint！(注意：转换后的PPT为图片格式)"
    except Exception as e:
        return False, f"转换失败: {str(e)}"

def pdf_to_jpg(input_path, output_dir):
    """
    将 PDF 的每一页导出为 JPG 图片
    """
    try:
        doc = fitz.open(input_path)
        base_name = os.path.splitext(os.path.basename(input_path))[0]

        for i in range(len(doc)):
            page = doc.load_page(i)
            # 缩放系数，提高图片清晰度 (144 dpi)
            zoom = 2.0
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            output_path = os.path.join(output_dir, f"{base_name}_page_{i+1}.jpg")
            pix.save(output_path)

        doc.close()
        return True, "PDF成功导出为JPG图片！"
    except Exception as e:
        return False, f"转换失败: {str(e)}"

def jpg_to_pdf(image_paths, output_path):
    """
    将多张 JPG/PNG 图片合并为一个 PDF
    """
    if not image_paths:
         return False, "没有选择图片。"
    try:
        # 打开第一张图片，并转为RGB
        images = []
        for path in image_paths:
             img = Image.open(path).convert('RGB')
             images.append(img)

        # 将后续图片追加到第一张图片的 PDF 格式保存中
        images[0].save(
            output_path,
            save_all=True,
            append_images=images[1:]
        )
        return True, "图片成功转换为PDF文件！"
    except Exception as e:
        return False, f"转换失败: {str(e)}"
